import csv
import os
import sys


def main() -> int:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception:
        print("Missing dependency: python-pptx. Install with: python -m pip install --user python-pptx")
        return 2

    root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    csv_path = os.path.join(root, "RAILWAY_PRESENTATION_SLIDES.csv")
    out_path = os.path.join(root, "RAILWAY_PRESENTATION_SLIDES.pptx")
    marker_path = os.path.join(root, "RAILWAY_PPT_GENERATION_MARKER.txt")

    if not os.path.exists(csv_path):
        print(f"CSV not found: {csv_path}")
        return 1

    prs = Presentation()

    with open(csv_path, "r", encoding="utf-8-sig", newline="") as f:
        reader = csv.DictReader(f)
        rows = list(reader)

    for row in rows:
        title = (row.get("Title") or "").strip()
        bullets = []
        for i in range(1, 7):
            v = row.get(f"Bullet {i}")
            if v is None:
                continue
            v = str(v).strip()
            if v:
                bullets.append(v)
        notes = (row.get("Speaker Notes") or "").strip()

        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title is not None:
            slide.shapes.title.text = title

        body = None
        for shp in slide.shapes:
            if shp.has_text_frame and shp != slide.shapes.title:
                body = shp
                break

        if body is not None:
            tf = body.text_frame
            tf.clear()
            tf.word_wrap = True

            if bullets:
                p0 = tf.paragraphs[0]
                p0.text = bullets[0]
                p0.level = 0
                for b in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0
            else:
                tf.text = ""

            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(20)

        if notes:
            ns = slide.notes_slide
            ntf = ns.notes_text_frame
            ntf.clear()
            ntf.text = notes

    prs.save(out_path)
    with open(marker_path, "w", encoding="utf-8") as mf:
        mf.write(out_path)
    print(f"Wrote: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

